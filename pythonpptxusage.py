from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_VERTICAL_ANCHOR, PP_PARAGRAPH_ALIGNMENT
from datetime import datetime
import os

# Opens a presentation called 'test' that already exists
print("\nusage: [name_of_powerpoint.pptx]")
name = input("Enter the name of a PowerPoint presentation to be generated: ")
prs = Presentation(name)
# prs = Presentation('test.pptx')
# prs.save('new.pptx')

prs = Presentation()

# -------- TITLE SLIDE -------- #
title_slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]

# Add text to layout
title.text = "TITLE HERE"

# Determines Date and Year text
datem = datetime.now().strftime('%B') # Format month
datey = datetime.now().strftime('%Y') # Format year
subtitle.text = datem + ' ' + datey
# Position text under title
subtitle.top = Inches(4)
subtitle.left = Inches(1)
subtitle.width = Inches(5)
subtitle.height = Inches(1)
# Change font size
subtitle.text_frame.paragraphs[0].font.size = Pt(36)


# -------- OVERVIEW -------- #
summary = prs.slide_layouts[6]
slide = prs.slides.add_slide(summary)
title = slide.shapes.title


# Add title
title.text = "TITLE"
# Position title
title.top = Inches(.02)
title.left = Inches(.5)
title.width = Inches(6)
title.height = Inches(.7)
# Change font size
title.text_frame.paragraphs[0].font.size = Pt(28)


# --Add Table--- #
rows, cols = 21, 8
x, y, cx, cy = Inches(.55), Inches(.6), Inches(9.2), Inches(1)  # position left, position top, width, height
shape = slide.shapes.add_table(rows, cols, x, y, cx, cy)  # rows, cols
table1 = shape.table

# Write column headers
overview_col_headers = ['TITLE', 'TITLE', 'TITLE', 'TITLE', 'TITLE', 'TITLE',
                    'TITLE', 'TITLE']
for i in range (8):
    table1.cell(0, i).text = overview_col_headers[i]

overview_row_headers = ['TITLE', 'TITLE', 'TITLE', 'TITLE', 'TITLE', 'TITLE', 'TITLE',
                        'TITLE', 'TITLE', 'TITLE', 'TITLE', 'TITLE', 'TITLE', 'TITLE', 'TITLE',
                        'TITLE', 'TITLE', 'TITLE', 'TITLE', 'TITLE']
for i in range (1, 21):
    table1.cell(i, 0).text = overview_row_headers[i-1]

# Format all cells
for r in range(rows):
    for c in range(cols):
        cell = table1.cell(r,c)
        # Text size
        cell.text_frame.paragraphs[0].font.size = Pt(9)
        # Margins
        cell.margin_left = Inches(.05)
        cell.margin_right = Inches(.05)
        cell.margin_top = Inches(.05)
        cell.margin_bottom = Inches(.05)
        # Fill color
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(0,158,47)

# Set column widths
table1.columns[0].width = Inches(1.2)
table1.columns[1].width = Inches(.7)
table1.columns[2].width = Inches(.6)
table1.columns[3].width = Inches(.6)
table1.columns[4].width = Inches(.8)
table1.columns[5].width = Inches(.8)
table1.columns[6].width = Inches(.7)
table1.columns[7].width = Inches(3.8)

# Format column headers
for r in range(1):
    for c in range(cols):
        cell = table1.cell(r,c)
        cell.vertical_anchor = MSO_VERTICAL_ANCHOR.BOTTOM
        cell.text_frame.paragraphs[0].alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(199, 199, 199)
        paragraph = cell.text_frame.paragraphs[0]
        paragraph.font.color.rgb = RGBColor(0, 0, 0)

# Format row headers
for r in range(1,rows):
    for c in range(1):
        cell = table1.cell(r,c)
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(255, 255, 255)
        cell.text_frame.paragraphs[0].font.bold = True

# Format Notes column color to 'white'
for r in range(1,rows):
    for c in range(7,cols):
        cell = table1.cell(r,c)
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(255, 255, 255)


# Add text boxes
txBox1 = slide.shapes.add_textbox(Inches(9.9), Inches(.9), Inches(3.2), Inches(6)) #left, top, width, height
summary = txBox1.text_frame

summary.text = "Summary Here"
summary.paragraphs[0].font.size = Pt(18)


# -------- SLIDE 1 -------- #
layout = prs.slide_layouts[6]  
slide = prs.slides.add_slide(layout)
title = slide.shapes.title

# Add text to layout
title.text = "TITLE HERE"
# Position title
title.top = Inches(.02)
title.left = Inches(.5)
title.width = Inches(6)
title.height = Inches(.7)
# Change font size
title.text_frame.paragraphs[0].font.size = Pt(28)


# Add 2 text boxes
txBox1 = slide.shapes.add_textbox(Inches(.55), Inches(.65), Inches(3.35), Inches(2.7))  # left, top, width, height
summary1 = txBox1.text_frame

summary1.text = "Summary Here"
summary1.paragraphs[0].font.size = Pt(12)
p = summary1.add_paragraph()
p.level = 0
p.text = "-Text Here"
summary1.paragraphs[1].font.size = Pt(12)
p2 = summary1.add_paragraph()
p2.level = 0
p2.text = "-Text Here"
summary1.paragraphs[2].font.size = Pt(12)
p3 = summary1.add_paragraph()
p3.level = 1
p3.text = "-Text Here"
summary1.paragraphs[3].font.size = Pt(12)

txBox2 = slide.shapes.add_textbox(Inches(.55), Inches(3.5), Inches(3.35), Inches(2.7))  # left, top, width, height
summary2 = txBox2.text_frame

summary2.text = "Summary Here"
summary2.paragraphs[0].font.size = Pt(12)
p = summary2.add_paragraph()
p.level = 0
p.text = "-Text Here"
summary2.paragraphs[1].font.size = Pt(12)
p2 = summary2.add_paragraph()
p2.level = 0
p2.text = "-Text Here"
summary2.paragraphs[2].font.size = Pt(12)
p3 = summary2.add_paragraph()
p3.level = 1
p3.text = "-Text Here"
summary2.paragraphs[3].font.size = Pt(12)

# Add table
tf = slide.shapes.add_textbox(Inches(4.89), Inches(0.5), Inches(2.8), Inches(.3))  # left, top, width, height
table_title = tf.text_frame
table_title.text = "TITLE HERE"
table_title.paragraphs[0].font.size = Pt(14)

cols, rows = 8, 9
x, y, cx, cy = Inches(4.97), Inches(.75), Inches(7), Inches(1)  # position left, position top, width, height
shape = slide.shapes.add_table(rows, cols, x, y, cx, cy)  # rows, cols
table1 = shape.table

# Write headers
col_headers = ['TITLE', 'TITLE', 'TITLE', 'TITLE', 'TITLE', 'TITLE', 'TITLE', 'TITLE']
row_headers = ['TITLE', 'TITLE', 'TITLE', 'TITLE', 'TITLE', 'TITLE', 'TITLE', 'TITLE']
values = ['####', '####', '####', '####', '####', '####', '####', '####']

for i in range(1, 9):
    table1.cell(0, i-1).text = col_headers[i-1]
    table1.cell(i, 0).text = row_headers[i-1]
    table1.cell(i, 1).text = values[i-1]

# Format all cells
for r in range(rows):
    for c in range(cols):
        cell = table1.cell(r, c)
        # Text size
        cell.text_frame.paragraphs[0].font.size = Pt(8)
        # Margins
        cell.margin_left = Inches(.05)
        cell.margin_right = Inches(.05)
        cell.margin_top = Inches(.05)
        cell.margin_bottom = Inches(.05)

# Set column widths
table1.columns[0].width = Inches(1.2)
table1.columns[1].width = Inches(.868)
for i in range(2, 8):
    table1.columns[i].width = Inches(.82)

# Format column headers
for r in range(1):
    for c in range(cols):
        cell = table1.cell(r, c)
        cell.vertical_anchor = MSO_VERTICAL_ANCHOR.BOTTOM
        cell.text_frame.paragraphs[0].alignment = PP_PARAGRAPH_ALIGNMENT.CENTER

# Format row headers
for r in range(1, rows):
    for c in range(1):
        cell = table1.cell(r, c)
        cell.text_frame.paragraphs[0].font.bold = True

# Center all data cells
for r in range(1, rows):
    for c in range(1, cols):
        cell = table1.cell(r, c)
        cell.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
        cell.text_frame.paragraphs[0].alignment = PP_PARAGRAPH_ALIGNMENT.CENTER


# -------- Save presentation -------- #
prs.save('new.pptx')
os.startfile('new.pptx')
